import os
import io
import re
import uuid
from typing import List, Optional, Literal, Dict, Any

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel

from database import db, create_document, get_documents
from schemas import Note, Quiz, Flashcard, Folder

# External libs for parsing
import requests

try:
    import docx
except Exception:
    docx = None

try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    import pptx
except Exception:
    pptx = None

app = FastAPI(title="Smart Notes Generator API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ----------------------------- Helpers -----------------------------
YOUTUBE_REGEX = r"(https?://)?(www\.)?(youtube\.com|youtu\.be)/.+"

def read_txt(file_bytes: bytes) -> str:
    return file_bytes.decode(errors="ignore")

def read_pdf(file_bytes: bytes) -> str:
    if PyPDF2 is None:
        raise HTTPException(500, "PDF parser not available")
    text = []
    reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
    for page in reader.pages:
        try:
            text.append(page.extract_text() or "")
        except Exception:
            continue
    return "\n".join(text)

def read_docx(file_bytes: bytes) -> str:
    if docx is None:
        raise HTTPException(500, "DOCX parser not available")
    tmp = io.BytesIO(file_bytes)
    d = docx.Document(tmp)
    return "\n".join(p.text for p in d.paragraphs)

def read_pptx(file_bytes: bytes) -> str:
    if pptx is None:
        raise HTTPException(500, "PPTX parser not available")
    prs = pptx.Presentation(io.BytesIO(file_bytes))
    slides_text = []
    for slide in prs.slides:
        chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
        slides_text.append("\n".join(chunks))
    return "\n\n".join(slides_text)

def fetch_url_text(url: str) -> str:
    # Very basic fetch; for YouTube, we can't fetch transcript here.
    if re.match(YOUTUBE_REGEX, url):
        # Placeholder: a simple instruction; real transcript handled by /transcribe
        return f"YouTube URL detected: {url}. Use transcription to extract audio text."
    try:
        r = requests.get(url, timeout=15)
        r.raise_for_status()
        return r.text
    except Exception as e:
        raise HTTPException(400, f"Failed to fetch URL: {e}")

# Simple keyword extraction
STOPWORDS = set("a an the and or if of in on to for with from by is are was were be been it this that as at into over about between among within without upon against due per via not but so than too very can will would could should might may have has had do does did".split())

def extract_keywords(text: str, limit: int = 15) -> List[str]:
    words = re.findall(r"[A-Za-z][A-Za-z\-]{2,}", text.lower())
    freq: Dict[str, int] = {}
    for w in words:
        if w in STOPWORDS:
            continue
        freq[w] = freq.get(w, 0) + 1
    return [w for w, _ in sorted(freq.items(), key=lambda kv: kv[1], reverse=True)[:limit]]

# Topic detection via simple heading/keyword heuristics

def detect_topics(text: str) -> List[str]:
    lines = [l.strip() for l in text.splitlines() if l.strip()]
    headings = [l for l in lines if len(l.split()) <= 8 and l.endswith(":")]
    if not headings:
        headings = [l for l in lines if len(l.split()) <= 6 and l.istitle()]
    topics = list(dict.fromkeys(headings))[:10]
    if not topics and text:
        topics = extract_keywords(text, 8)
    return topics

# Pseudo AI generation: deterministic template-based generator (no external API)
# We avoid calling closed models; produces structured notes that react to options.

class GenerateOptions(BaseModel):
    note_type: Literal["bullet","detailed","summary","exam","timestamp","mindmap"]
    tone: Literal["simple","professional","technical"]
    length: Literal["short","medium","detailed"]
    highlight_key_points: bool = True
    include_examples: bool = False
    extract_definitions: bool = True
    include_formulas: bool = False


def style_text(text: str, tone: str) -> str:
    if tone == "simple":
        return text
    if tone == "professional":
        return re.sub(r"\b(\w{4,})\b", lambda m: m.group(1).capitalize(), text)
    if tone == "technical":
        return text + "\n\nTechnical Notes: Terms aligned with domain-specific accuracy."
    return text


def generate_notes_from_text(text: str, options: GenerateOptions) -> str:
    keywords = extract_keywords(text)
    topics = detect_topics(text)
    chunk_factor = {"short": 0.4, "medium": 0.8, "detailed": 1.2}[options.length]
    base_len = max(10, min(60, int(len(text.split()) * 0.03 * chunk_factor)))

    def bullets():
        pts = topics or ["Key Points"]
        lines = [f"# Smart Notes", "", f"## Topics: {', '.join(topics[:6]) if topics else 'General'}", ""]
        lines.append("### Bullet Points")
        sample = keywords[: min(10, base_len)] or ["Overview", "Concepts"]
        for w in sample:
            lines.append(f"- {w.title()} — concise explanation relevant to the lecture context.")
        if options.include_examples:
            lines.append("\n#### Examples")
            lines.append("- Example 1: Practical case illustrating the concept.")
        if options.extract_definitions:
            lines.append("\n#### Definitions")
            for k in sample[:5]:
                lines.append(f"- {k.title()}: Clear definition explained simply.")
        return "\n".join(lines)

    def detailed():
        sections = topics[:5] or ["Overview", "Details"]
        out = ["# Smart Notes (Detailed)", ""]
        for s in sections:
            out.append(f"## {s}")
            out.append("- Background and context with clarity.")
            out.append("- Important relationships or equations.")
            if options.include_examples:
                out.append("- Example: A concrete scenario to solidify understanding.")
            if options.include_formulas:
                out.append("- Formula: y = f(x) with assumptions and units.")
            out.append("")
        out.append("## Key Takeaways")
        for k in keywords[:10]:
            out.append(f"- {k.title()}")
        return "\n".join(out)

    def summary():
        out = ["# Summary Notes", ""]
        out.append("## Core Idea")
        out.append("A concise overview eliminating filler, focusing on what matters.")
        out.append("\n## Highlights")
        for k in keywords[:8]:
            out.append(f"- {k.title()}")
        return "\n".join(out)

    def exam():
        out = ["# Exam-Ready Notes", ""]
        out.append("## Must-Know Points")
        for k in keywords[:12]:
            out.append(f"- {k.title()} — Definition and quick recall cue.")
        out.append("\n## Common Pitfalls")
        out.append("- Confusing similar terms, forgetting boundary cases.")
        out.append("\n## Quick Formulas")
        if options.include_formulas:
            out.append("- Key Equations with units and conditions.")
        return "\n".join(out)

    def timestamp():
        out = ["# Timestamp-Based Notes", ""]
        out.append("## Segments")
        for i, k in enumerate(keywords[:10]):
            out.append(f"- [0{i}:00] {k.title()} — Key discussion summary.")
        return "\n".join(out)

    def mindmap():
        out = ["# Mindmap Outline", ""]
        root = topics[0] if topics else "Topic"
        out.append(f"- {root}")
        for k in keywords[:8]:
            out.append(f"  - {k.title()}")
            out.append(f"    - Detail about {k}.")
        return "\n".join(out)

    mapping = {
        "bullet": bullets,
        "detailed": detailed,
        "summary": summary,
        "exam": exam,
        "timestamp": timestamp,
        "mindmap": mindmap,
    }

    result = mapping[options.note_type]()
    if options.highlight_key_points:
        result = result.replace("- ", "- ❖ ")
    result = style_text(result, options.tone)
    return result

# ----------------------------- Models -----------------------------
class GenerateRequest(BaseModel):
    text: str
    options: GenerateOptions
    title: Optional[str] = None
    source_type: Literal["text","file","url"] = "text"

class SaveNoteRequest(BaseModel):
    title: str
    content: str
    folder_id: Optional[str] = None
    options: Dict[str, Any] = {}
    transcript: Optional[str] = None

# ----------------------------- Routes -----------------------------
@app.get("/")
async def root():
    return {"name": "Smart Notes Generator API", "status": "ok"}

@app.get("/schema")
async def get_schema():
    return {
        "collections": ["folder", "note", "quiz", "flashcard", "setting"],
    }

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    ext = (file.filename.split(".")[-1] or "").lower()
    content_bytes = await file.read()
    if ext in ["txt", "md"]:
        text = read_txt(content_bytes)
    elif ext == "pdf":
        text = read_pdf(content_bytes)
    elif ext in ["docx"]:
        text = read_docx(content_bytes)
    elif ext in ["ppt", "pptx"]:
        text = read_pptx(content_bytes)
    elif ext in ["mp3", "wav", "aac", "m4a", "mp4", "mkv", "mov"]:
        # For media, respond with indicator; real transcript via /transcribe
        text = "[MEDIA FILE UPLOADED] Use /transcribe to extract transcript."
    else:
        raise HTTPException(400, f"Unsupported file type: .{ext}")

    return {"filename": file.filename, "text": text[:100000]}

@app.post("/extract-url")
async def extract_url(url: str = Form(...)):
    txt = fetch_url_text(url)
    return {"url": url, "text": txt[:100000]}

@app.post("/generate")
async def generate(req: GenerateRequest):
    if not req.text or len(req.text.strip()) < 3:
        raise HTTPException(400, "Provide text or transcript for generation")
    notes = generate_notes_from_text(req.text, req.options)
    title = req.title or (req.text.splitlines()[0][:60] if req.text else "Smart Notes")
    return {"title": title, "content": notes, "keywords": extract_keywords(req.text), "topics": detect_topics(req.text)}

@app.post("/save")
async def save_note(req: SaveNoteRequest):
    data = Note(
        title=req.title,
        content=req.content,
        source_type="text",
        source_name=req.options.get("source_name"),
        options=req.options,
        folder_id=req.folder_id,
        keywords=extract_keywords(req.content),
        topics=detect_topics(req.content),
        transcript=req.transcript,
    )
    inserted_id = create_document("note", data)
    return {"id": inserted_id}

@app.get("/notes")
async def list_notes():
    docs = get_documents("note", {}, limit=100)
    # Convert ObjectId to str
    for d in docs:
        d["_id"] = str(d["_id"])  # type: ignore
    return {"items": docs}

@app.get("/notes/{note_id}")
async def get_note(note_id: str):
    from bson import ObjectId
    try:
        doc = db["note"].find_one({"_id": ObjectId(note_id)})
        if not doc:
            raise HTTPException(404, "Note not found")
        doc["_id"] = str(doc["_id"])  # type: ignore
        return doc
    except Exception:
        raise HTTPException(400, "Invalid note id")

@app.post("/folders")
async def create_folder(name: str = Form(...), color: Optional[str] = Form(None)):
    data = Folder(name=name, color=color)
    fid = create_document("folder", data)
    return {"id": fid}

@app.get("/folders")
async def list_folders():
    docs = get_documents("folder", {}, limit=100)
    for d in docs:
        d["_id"] = str(d["_id"])  # type: ignore
    return {"items": docs}

# Naive transcription placeholder using text upload (no external Whisper)
@app.post("/transcribe")
async def transcribe(file: UploadFile = File(...)):
    ext = (file.filename.split(".")[-1] or "").lower()
    content_bytes = await file.read()
    if ext in ["mp3", "wav", "aac", "m4a", "mp4", "mkv", "mov"]:
        # We cannot process audio/video in this environment to real text reliably.
        # For MVP, we return a stub stating media content available and ask user to provide rough text if needed.
        transcript = "Transcription placeholder: Please provide approximate text or notes; real deployment should integrate Whisper API."
    else:
        raise HTTPException(400, "File must be audio/video for transcription")
    return {"transcript": transcript}

# Simple quiz generator
class QuizRequest(BaseModel):
    note_id: Optional[str] = None
    content: Optional[str] = None
    num_questions: int = 5

@app.post("/quiz")
async def generate_quiz(req: QuizRequest):
    text = req.content
    if not text and req.note_id:
        from bson import ObjectId
        doc = db["note"].find_one({"_id": ObjectId(req.note_id)})
        if not doc:
            raise HTTPException(404, "Note not found")
        text = doc.get("content", "")
    if not text:
        raise HTTPException(400, "No content to generate quiz")
    kws = extract_keywords(text, limit=req.num_questions * 2)
    qs = []
    for i, k in enumerate(kws[: req.num_questions]):
        opts = [k.title(), f"Not {k}", f"Unrelated {i}", f"Opposite {i}"]
        qs.append({"q": f"What best describes {k}?", "options": opts, "answer": 0})
    return {"title": "Auto MCQ Quiz", "questions": qs}

# Flashcards
class FlashcardRequest(BaseModel):
    content: str
    num_cards: int = 10

@app.post("/flashcards")
async def generate_flashcards(req: FlashcardRequest):
    kws = extract_keywords(req.content, limit=req.num_cards)
    cards = [{"front": k.title(), "back": f"Definition and usage of {k}."} for k in kws]
    return {"cards": cards}

# Export endpoints - return text/markdown download; for PDF/DOCX in frontend
@app.post("/export/txt")
async def export_txt(title: str = Form(...), content: str = Form(...)):
    filename = f"{re.sub(r'[^A-Za-z0-9_-]+', '_', title) or 'notes'}.txt"
    def iterfile():
        yield content
    headers = {"Content-Disposition": f"attachment; filename={filename}"}
    return StreamingResponse(iterfile(), media_type="text/plain", headers=headers)

@app.post("/export/md")
async def export_md(title: str = Form(...), content: str = Form(...)):
    filename = f"{re.sub(r'[^A-Za-z0-9_-]+', '_', title) or 'notes'}.md"
    def iterfile():
        yield content
    headers = {"Content-Disposition": f"attachment; filename={filename}"}
    return StreamingResponse(iterfile(), media_type="text/markdown", headers=headers)

# Health/database test
@app.get("/test")
def test_database():
    response = {
        "backend": "✅ Running",
        "database": "❌ Not Available",
        "database_url": None,
        "database_name": None,
        "connection_status": "Not Connected",
        "collections": []
    }
    try:
        if db is not None:
            response["database"] = "✅ Available"
            response["database_url"] = "✅ Configured"
            response["database_name"] = db.name if hasattr(db, 'name') else "✅ Connected"
            try:
                collections = db.list_collection_names()
                response["collections"] = collections[:10]
                response["database"] = "✅ Connected & Working"
            except Exception as e:
                response["database"] = f"⚠️  Connected but Error: {str(e)[:50]}"
        else:
            response["database"] = "⚠️  Available but not initialized"
    except Exception as e:
        response["database"] = f"❌ Error: {str(e)[:50]}"

    response["database_url"] = "✅ Set" if os.getenv("DATABASE_URL") else "❌ Not Set"
    response["database_name"] = "✅ Set" if os.getenv("DATABASE_NAME") else "❌ Not Set"
    return response

if __name__ == "__main__":
    import uvicorn
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
